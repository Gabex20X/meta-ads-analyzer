import streamlit as st
import pandas as pd
import google.generativeai as genai
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import json
import io
import os
from datetime import datetime

# ── [REGRA 2] Importações para exportação PDF e PPTX ─────────────────────────
from fpdf import FPDF                              # fpdf2 — suporte UTF-8 nativo
from pptx import Presentation                      # python-pptx
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
# ─────────────────────────────────────────────────────────────────────────────

# ── Page config ──────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Meta Ads AI Analyzer",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Custom CSS ────────────────────────────────────────────────────────────────
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Space+Mono:wght@400;700&family=Syne:wght@400;600;800&display=swap');

/* ── Root & Reset ── */
:root {
    --bg:        #0a0a0f;
    --surface:   #12121a;
    --border:    #2a2a3a;
    --accent:    #7c6aff;
    --accent2:   #ff6a9b;
    --text:      #e8e6ff;
    --muted:     #6b6880;
    --success:   #4fffb0;
    --warning:   #ffcc4f;
    --danger:    #ff4f6a;
}

html, body, [class*="css"] {
    font-family: 'Syne', sans-serif;
    background: var(--bg) !important;
    color: var(--text) !important;
}

/* ── App bg ── */
.stApp {
    background: var(--bg) !important;
}

/* ── Sidebar ── */
[data-testid="stSidebar"] {
    background: var(--surface) !important;
    border-right: 1px solid var(--border) !important;
}

/* ── Header ── */
.hero-header {
    padding: 2rem 0 1.5rem;
    text-align: center;
}
.hero-header h1 {
    font-size: 2.8rem;
    font-weight: 800;
    letter-spacing: -1px;
    background: linear-gradient(135deg, var(--accent) 0%, var(--accent2) 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    margin-bottom: 0.3rem;
}
.hero-header p {
    color: var(--muted);
    font-family: 'Space Mono', monospace;
    font-size: 0.8rem;
    letter-spacing: 0.1em;
    text-transform: uppercase;
}

/* ── KPI Cards ── */
.kpi-grid {
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(180px, 1fr));
    gap: 1rem;
    margin: 1.5rem 0;
}
.kpi-card {
    background: var(--surface);
    border: 1px solid var(--border);
    border-radius: 12px;
    padding: 1.2rem 1.5rem;
    position: relative;
    overflow: hidden;
    transition: border-color 0.2s;
}
.kpi-card::before {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0;
    height: 3px;
}
.kpi-card.purple::before { background: var(--accent); }
.kpi-card.pink::before   { background: var(--accent2); }
.kpi-card.green::before  { background: var(--success); }
.kpi-card.yellow::before { background: var(--warning); }
.kpi-card.red::before    { background: var(--danger); }

.kpi-label {
    font-family: 'Space Mono', monospace;
    font-size: 0.65rem;
    text-transform: uppercase;
    letter-spacing: 0.12em;
    color: var(--muted);
    margin-bottom: 0.4rem;
}
.kpi-value {
    font-size: 1.7rem;
    font-weight: 800;
    line-height: 1;
    color: var(--text);
}
.kpi-sub {
    font-size: 0.72rem;
    color: var(--muted);
    margin-top: 0.3rem;
}

/* ── Section titles ── */
.section-title {
    font-size: 1rem;
    font-weight: 600;
    color: var(--text);
    letter-spacing: 0.03em;
    margin: 2rem 0 0.8rem;
    display: flex;
    align-items: center;
    gap: 0.5rem;
}
.section-title::after {
    content: '';
    flex: 1;
    height: 1px;
    background: var(--border);
}

/* ── AI Analysis box ── */
.ai-box {
    background: linear-gradient(135deg, #12121a 60%, #1a1228 100%);
    border: 1px solid var(--accent);
    border-radius: 14px;
    padding: 1.8rem 2rem;
    position: relative;
    overflow: hidden;
}
.ai-box::before {
    content: '';
    position: absolute;
    top: -60px; right: -60px;
    width: 180px; height: 180px;
    background: radial-gradient(circle, rgba(124,106,255,0.15) 0%, transparent 70%);
    pointer-events: none;
}
.ai-tag {
    display: inline-block;
    background: rgba(124,106,255,0.15);
    border: 1px solid rgba(124,106,255,0.4);
    color: var(--accent);
    font-family: 'Space Mono', monospace;
    font-size: 0.65rem;
    letter-spacing: 0.1em;
    text-transform: uppercase;
    padding: 0.25rem 0.7rem;
    border-radius: 20px;
    margin-bottom: 1rem;
}
.ai-content {
    color: #c8c6e8;
    line-height: 1.75;
    font-size: 0.95rem;
    white-space: pre-wrap;
}

/* ── Dataframe ── */
[data-testid="stDataFrame"] {
    border-radius: 10px;
    overflow: hidden;
    border: 1px solid var(--border) !important;
}

/* ── Buttons ── */
.stButton > button {
    background: linear-gradient(135deg, var(--accent), var(--accent2)) !important;
    color: white !important;
    border: none !important;
    border-radius: 8px !important;
    font-family: 'Syne', sans-serif !important;
    font-weight: 600 !important;
    padding: 0.6rem 1.8rem !important;
    transition: opacity 0.2s, transform 0.1s !important;
}
.stButton > button:hover {
    opacity: 0.88 !important;
    transform: translateY(-1px) !important;
}

/* ── File uploader ── */
[data-testid="stFileUploader"] {
    border: 2px dashed var(--border) !important;
    border-radius: 12px !important;
    background: var(--surface) !important;
}

/* ── Text inputs ── */
.stTextInput input, .stSelectbox > div > div {
    background: var(--surface) !important;
    border: 1px solid var(--border) !important;
    color: var(--text) !important;
    border-radius: 8px !important;
}

/* ── Spinner text ── */
.stSpinner > div {
    color: var(--accent) !important;
}

/* ── Alert-style info ── */
.stAlert {
    background: var(--surface) !important;
    border: 1px solid var(--border) !important;
    border-radius: 10px !important;
}
</style>
""", unsafe_allow_html=True)


# ── Helpers ──────────────────────────────────────────────────────────────────

REQUIRED_COLS = [
    "campaign_name", "impressions", "clicks", "spend", "reach"
]

OPTIONAL_COLS = [
    "ctr", "cpc", "cpm", "frequency", "conversions",
    "conversion_values", "roas", "ad_name", "adset_name",
    "result_rate", "cost_per_result"
]

def fmt_currency(val: float, symbol: str = "R$") -> str:
    return f"{symbol} {val:,.2f}"

def fmt_number(val: float) -> str:
    if val >= 1_000_000:
        return f"{val/1_000_000:.2f}M"
    if val >= 1_000:
        return f"{val/1_000:.1f}K"
    return f"{val:,.0f}"

def fmt_pct(val: float) -> str:
    return f"{val:.2f}%"

def load_csv(file) -> pd.DataFrame:
    """Try multiple encodings to load the CSV."""
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            file.seek(0)
            df = pd.read_csv(file, encoding=enc)
            df.columns = [c.strip().lower().replace(" ", "_") for c in df.columns]
            return df
        except Exception:
            continue
    raise ValueError("Não foi possível decodificar o arquivo CSV.")

# Regra 2: cada plataforma tem nomes de colunas distintos. Esta função mapeia
# tudo para o padrão interno único usado pelos gráficos e cálculos de KPI.
# Adicionar novas plataformas no futuro = apenas um novo bloco de mapeamento.

# Mapeamento Google Ads → padrão interno
# Chave: nome da coluna do CSV (lowercase+underscore), Valor: nome interno
GOOGLE_ADS_COLUMN_MAP: dict[str, str] = {
    # Campanha
    "campaign":                "campaign_name",
    "campanha":                "campaign_name",
    # Impressões
    "impr.":                   "impressions",
    "impressões":              "impressions",
    "impressions":             "impressions",
    # Cliques
    "clicks":                  "clicks",
    "cliques":                 "clicks",
    # Investimento
    "cost":                    "spend",
    "custo":                   "spend",
    "cost_(brl)":              "spend",
    "custo_(brl)":             "spend",
    "cost_(usd)":              "spend",
    # CTR
    "ctr":                     "ctr",
    # CPC médio
    "avg._cpc":                "cpc",
    "cpc_médio":               "cpc",
    "avg._cpc_(brl)":          "cpc",
    # CPM
    "avg._cpm":                "cpm",
    "cpm_médio":               "cpm",
    # Conversões
    "conversions":             "conversions",
    "conversões":              "conversions",
    # Valor das conversões (para ROAS)
    "conv._value":             "conversion_values",
    "valor_conv.":             "conversion_values",
    "all_conv._value":         "conversion_values",
    # Alcance (Google não exporta reach diretamente)
    "reach":                   "reach",
    # Ad / grupo de anúncios
    "ad_group":                "adset_name",
    "grupo_de_anúncios":       "adset_name",
    "ad":                      "ad_name",
    "anúncio":                 "ad_name",
    # Quota de impressão (específico Google)
    "search_impr._share":      "impression_share",
    # Data (Google exporta como "Day", "Week", "Month")
    "day":                     "date",
    "dia":                     "date",
    "week":                    "week",
    "month":                   "month",
}

# Meta Ads já usa nomes padrão; mapeamento identidade apenas para os aliases
META_ADS_COLUMN_MAP: dict[str, str] = {
    # Aliases PT-BR que podem vir no export localizado
    "campanha":            "campaign_name",
    "investimento":        "spend",
    "impressões":          "impressions",
    "cliques":             "clicks",
    "alcance":             "reach",
    "conversões":          "conversions",
    "frequência":          "frequency",
    "nome_do_anúncio":     "ad_name",
    "conjunto_de_anúncios":"adset_name",
}


def normalize_columns(df: pd.DataFrame, platform: str) -> pd.DataFrame:
    """
    Aplica o mapeamento de colunas da plataforma selecionada para o padrão
    interno. Colunas não mapeadas são mantidas intactas.
    """
    col_map = GOOGLE_ADS_COLUMN_MAP if platform == "Google Ads" else META_ADS_COLUMN_MAP
    # Normaliza os nomes atuais do df para lowercase+underscore antes de comparar
    rename = {col: col_map[col] for col in df.columns if col in col_map}
    df = df.rename(columns=rename)
    return df

# Regra 3: o Google Ads exporta uma linha de título do relatório no topo e
# uma linha de "Total" no final. Ambas precisam ser removidas antes do parse
# para não contaminar os dados numéricos.

def clean_google_ads_csv(file) -> pd.DataFrame:
    """
    Lê o CSV do Google Ads ignorando linhas de cabeçalho extras (relatório
    gerado pelo Google tem 2 linhas antes do header real) e remove a linha
    de totais no final.
    """
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            file.seek(0)
            raw = file.read().decode(enc)
            lines = raw.splitlines()

            # Encontra a linha do header real: é a primeira linha que
            # contenha "Campaign" ou "Campanha" (case-insensitive).
            header_idx = 0
            for i, line in enumerate(lines):
                lower = line.lower()
                if "campaign" in lower or "campanha" in lower or "impr" in lower:
                    header_idx = i
                    break

            # Reconstrói o CSV a partir do header real
            clean = "\n".join(lines[header_idx:])
            import io as _io
            df = pd.read_csv(_io.StringIO(clean))
            df.columns = [c.strip().lower().replace(" ", "_") for c in df.columns]

            # Remove linhas de total: onde a coluna de campanha é "total"
            # (Google Ads exporta "Total" em PT e EN)
            camp_col = next(
                (c for c in df.columns if c in ("campaign", "campanha", "campaign_name")),
                None,
            )
            if camp_col:
                df = df[~df[camp_col].astype(str).str.strip().str.lower().isin(
                    ["total", "totals", "grand total", "total geral"]
                )]

            # Remove linhas completamente vazias
            df = df.dropna(how="all").reset_index(drop=True)
            return df

        except Exception:
            continue

    raise ValueError("Não foi possível decodificar o CSV do Google Ads.")

def preprocess(df: pd.DataFrame) -> pd.DataFrame:
    """Clean and derive metrics."""
    # Numeric coercion
    numeric_cols = [
        "impressions", "clicks", "spend", "reach", "ctr", "cpc",
        "cpm", "frequency", "conversions", "conversion_values",
        "roas", "result_rate", "cost_per_result"
    ]
    for col in numeric_cols:
        if col in df.columns:
            df[col] = (
                df[col]
                .astype(str)
                .str.replace(r"[R$\s%,]", "", regex=True)
                .str.replace(",", ".", regex=False)
                .pipe(pd.to_numeric, errors="coerce")
            )

    # Derive missing metrics
    if "ctr" not in df.columns and {"clicks", "impressions"}.issubset(df.columns):
        df["ctr"] = df["clicks"] / df["impressions"].replace(0, float("nan")) * 100

    if "cpc" not in df.columns and {"spend", "clicks"}.issubset(df.columns):
        df["cpc"] = df["spend"] / df["clicks"].replace(0, float("nan"))

    if "cpm" not in df.columns and {"spend", "impressions"}.issubset(df.columns):
        df["cpm"] = df["spend"] / df["impressions"].replace(0, float("nan")) * 1000

    if "roas" not in df.columns and {"conversion_values", "spend"}.issubset(df.columns):
        df["roas"] = df["conversion_values"] / df["spend"].replace(0, float("nan"))

    return df

def build_summary(df: pd.DataFrame) -> dict:
    """Build a compact summary dict for the AI prompt."""
    s: dict = {}

    def safe_sum(col):
        return float(df[col].sum()) if col in df.columns else None

    def safe_mean(col):
        return float(df[col].mean()) if col in df.columns else None

    s["total_campaigns"] = int(df["campaign_name"].nunique()) if "campaign_name" in df.columns else len(df)
    s["total_ads"] = int(df["ad_name"].nunique()) if "ad_name" in df.columns else None
    s["total_spend"] = safe_sum("spend")
    s["total_impressions"] = safe_sum("impressions")
    s["total_clicks"] = safe_sum("clicks")
    s["total_conversions"] = safe_sum("conversions")
    s["total_revenue"] = safe_sum("conversion_values")
    s["avg_ctr"] = safe_mean("ctr")
    s["avg_cpc"] = safe_mean("cpc")
    s["avg_cpm"] = safe_mean("cpm")
    s["avg_roas"] = safe_mean("roas")
    s["avg_frequency"] = safe_mean("frequency")

    # Top / bottom campaigns by ROAS or spend
    if "campaign_name" in df.columns and "spend" in df.columns:
        grp = df.groupby("campaign_name").agg(
            spend=("spend", "sum"),
            clicks=("clicks", "sum") if "clicks" in df.columns else ("spend", "count"),
        ).reset_index().sort_values("spend", ascending=False)
        s["top_campaigns_by_spend"] = grp.head(5).to_dict(orient="records")

    return {k: v for k, v in s.items() if v is not None}

# Regra 4: a plataforma é injetada no prompt para que o Gemini assuma o papel
# correto de especialista e use a nomenclatura certa (ex: "Google Ads Manager"
# vs "Meta Ads Manager" nas recomendações de implementação).

def build_prompt(summary: dict, user_goal: str, currency: str, platform: str) -> str:
    summary_json = json.dumps(summary, ensure_ascii=False, indent=2)

    # Textos que variam por plataforma
    platform_role = (
        "Meta Ads (Facebook/Instagram)"
        if platform == "Meta Ads"
        else "Google Ads (Search, Display, Performance Max)"
    )
    platform_manager = (
        "Meta Ads Manager"
        if platform == "Meta Ads"
        else "Google Ads Manager / Editor"
    )

    return f"""Você é um especialista sênior em performance de mídia paga com foco em {platform_role}. Analise os dados abaixo e forneça uma análise aprofundada e acionável em português.

## PLATAFORMA
{platform} — use terminologia, benchmarks e recomendações específicos desta plataforma.

## DADOS DOS ANÚNCIOS
```json
{summary_json}
```

## OBJETIVO DO ANUNCIANTE
{user_goal}

## MOEDA
{currency}

## ESTRUTURA DA ANÁLISE (responda nesta ordem exata):

### 🔍 Diagnóstico Geral
Avalie o desempenho geral da conta: volume de investimento, alcance, engajamento e eficiência de custo. Seja direto e específico com os números.

### 🏆 Pontos Fortes
Liste de 2 a 4 pontos positivos concretos (com métricas), explicando por que são bons resultados para {platform}.

### ⚠️ Alertas & Problemas
Identifique de 2 a 4 problemas críticos ou métricas abaixo do ideal, com benchmarks do setor para {platform} como comparação.

### 🎯 Recomendações Prioritárias
Forneça de 3 a 5 ações específicas e priorizadas (da mais urgente para a menos urgente). Cada ação deve incluir:
- O que fazer
- Por que fazer (impacto esperado)
- Como implementar (passos concretos no {platform_manager})

### 📈 Projeção de Oportunidade
Com base nos dados, estime o potencial de melhoria se as principais recomendações forem implementadas (ex.: redução de CPC, aumento de ROAS, etc.).

Seja direto, técnico e evite generalizações. Use os números reais dos dados fornecidos.
"""


def run_gemini_analysis(api_key: str, prompt: str, model_name: str) -> str:
    genai.configure(api_key=api_key)
    modelo_oficial = "gemini-2.5-flash"
    model = genai.GenerativeModel(modelo_oficial)
    response = model.generate_content(prompt)
    return response.text


# ── [REGRA 1] Stop Loss — lógica de diagnóstico de risco ─────────────────────
# Varre o DataFrame por campanha e retorna uma lista de alertas estruturados.
# Cada alerta tem: tipo ("error" | "warning"), campanha e mensagem.
# A função é agnóstica à plataforma — usa apenas colunas do padrão interno.

STOP_LOSS_SPEND_THRESHOLD = 100   # Gasto mínimo (moeda local) para alertar verba queimada
STOP_LOSS_CTR_MIN         = 0.50  # CTR mínimo aceitável em %
STOP_LOSS_ROAS_MIN        = 1.0   # ROAS mínimo para campanhas com conversão

def run_stop_loss(df: pd.DataFrame, sym: str) -> list[dict]:
    """
    Percorre o DataFrame agrupado por campanha e retorna alertas de risco.
    Retorna lista de dicts: {level: 'error'|'warning', campaign: str, msg: str}
    """
    alerts = []

    # Se não tiver campanha, analisa o DataFrame inteiro como uma "conta"
    group_col = "campaign_name" if "campaign_name" in df.columns else None
    groups = df.groupby(group_col) if group_col else [("Conta geral", df)]

    for camp_name, group in groups:
        spend       = group["spend"].sum()       if "spend"       in group.columns else None
        conversions = group["conversions"].sum()  if "conversions" in group.columns else None
        ctr         = group["ctr"].mean()         if "ctr"         in group.columns else None
        roas        = group["roas"].mean()        if "roas"        in group.columns else None

        # ── Alerta 1: Queima de Verba ───────────────────────────────────────
        # Gastou acima do threshold mas não gerou nenhuma conversão
        if spend is not None and conversions is not None:
            if spend > STOP_LOSS_SPEND_THRESHOLD and conversions == 0:
                alerts.append({
                    "level":    "error",
                    "campaign": camp_name,
                    "msg": (
                        f"🔥 **Queima de Verba** — Investimento de {sym} {spend:,.2f} "
                        f"sem nenhuma conversão registrada. Pausar ou revisar urgentemente."
                    ),
                })

        # ── Alerta 2: Baixa Relevância (CTR) ────────────────────────────────
        # CTR abaixo de 0.50% indica criativos ou segmentação fracos
        if ctr is not None and ctr < STOP_LOSS_CTR_MIN:
            alerts.append({
                "level":    "warning",
                "campaign": camp_name,
                "msg": (
                    f"📉 **Baixa Relevância** — CTR médio de {ctr:.2f}% está abaixo do "
                    f"mínimo recomendado de {STOP_LOSS_CTR_MIN}%. "
                    f"Revisar criativos e segmentação."
                ),
            })

        # ── Alerta 3: Prejuízo (ROAS < 1) ───────────────────────────────────
        # Só alertar se houve conversão para não confundir campanhas de topo de funil
        if roas is not None and conversions is not None and conversions > 0:
            if roas < STOP_LOSS_ROAS_MIN:
                alerts.append({
                    "level":    "error",
                    "campaign": camp_name,
                    "msg": (
                        f"🚨 **Prejuízo** — ROAS de {roas:.2f}x está abaixo de 1.0. "
                        f"A campanha está gerando menos receita do que consome. "
                        f"Ação imediata necessária."
                    ),
                })

    return alerts
# ── [fim REGRA 1] ─────────────────────────────────────────────────────────────


# ── [REGRA 2] Geradores de relatório — PDF e PPTX ────────────────────────────

def _kpi_lines(summary: dict, sym: str) -> list[str]:
    """Formata os KPIs do summary como linhas de texto simples para relatórios."""
    mapping = {
        "total_campaigns":   ("Campanhas",      lambda v: str(int(v))),
        "total_spend":       ("Investimento",   lambda v: f"{sym} {v:,.2f}"),
        "total_impressions": ("Impressões",     lambda v: f"{v:,.0f}"),
        "total_clicks":      ("Cliques",        lambda v: f"{v:,.0f}"),
        "avg_ctr":           ("CTR Médio",      lambda v: f"{v:.2f}%"),
        "avg_cpc":           ("CPC Médio",      lambda v: f"{sym} {v:,.2f}"),
        "avg_roas":          ("ROAS Médio",     lambda v: f"{v:.2f}x"),
        "total_conversions": ("Conversões",     lambda v: f"{v:,.0f}"),
        "total_revenue":     ("Receita Total",  lambda v: f"{sym} {v:,.2f}"),
    }
    lines = []
    for key, (label, fmt) in mapping.items():
        if key in summary and summary[key] is not None:
            lines.append(f"{label}: {fmt(summary[key])}")
    return lines


def generate_pdf(
    platform: str,
    summary: dict,
    alerts: list[dict],
    analysis_text: str,
    sym: str,
) -> bytes:
    """
    Gera um PDF completo com:
      - Título e metadados
      - Resumo dos KPIs
      - Alertas do Stop Loss
      - Análise estratégica completa do Gemini
    Retorna os bytes do PDF prontos para st.download_button.
    """

    class UTF8PDF(FPDF):
        """FPDF com suporte a caracteres UTF-8 via fonte DejaVu embutida."""
        def header(self):
            # Barra de topo colorida
            self.set_fill_color(124, 106, 255)   # roxo #7c6aff
            self.rect(0, 0, 210, 8, "F")
            self.ln(10)

        def footer(self):
            self.set_y(-15)
            self.set_font("DejaVu", size=8)
            self.set_text_color(150, 150, 170)
            self.cell(0, 10, f"Marketing Digital Hub  ·  {platform}  ·  Página {self.page_no()}", align="C")

    pdf = UTF8PDF()
    # fpdf2: adiciona fonte com suporte Unicode (DejaVu está embutida)
    pdf.add_font("DejaVu", fname="DejaVuSans.ttf")
    pdf.add_font("DejaVu", style="B", fname="DejaVuSans-Bold.ttf")
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_page()

    # ── Título ──
    pdf.set_font("DejaVu", style="B", size=22)
    pdf.set_text_color(124, 106, 255)
    pdf.cell(0, 12, "Marketing Digital Hub", ln=True, align="C")

    pdf.set_font("DejaVu", size=11)
    pdf.set_text_color(107, 104, 128)
    pdf.cell(0, 7, f"Plataforma: {platform}  ·  Gerado em: {datetime.now().strftime('%d/%m/%Y %H:%M')}", ln=True, align="C")
    pdf.ln(6)

    # ── Linha separadora ──
    pdf.set_draw_color(42, 42, 58)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(6)

    # ── KPIs ──
    pdf.set_font("DejaVu", style="B", size=13)
    pdf.set_text_color(232, 230, 255)
    pdf.cell(0, 8, "Resumo de KPIs", ln=True)
    pdf.ln(2)

    kpi_lines = _kpi_lines(summary, sym)
    pdf.set_font("DejaVu", size=10)
    pdf.set_text_color(200, 198, 232)
    for i, line in enumerate(kpi_lines):
        # Duas colunas
        if i % 2 == 0:
            pdf.cell(95, 7, line)
        else:
            pdf.cell(95, 7, line, ln=True)
    if len(kpi_lines) % 2 != 0:
        pdf.ln(7)
    pdf.ln(5)

    # ── Stop Loss ──
    pdf.set_font("DejaVu", style="B", size=13)
    pdf.set_text_color(232, 230, 255)
    pdf.cell(0, 8, "Diagnostico de Risco (Stop Loss)", ln=True)
    pdf.ln(2)

    pdf.set_font("DejaVu", size=10)
    if not alerts:
        pdf.set_text_color(79, 255, 176)   # verde #4fffb0
        pdf.multi_cell(0, 6, "Todas as campanhas estao dentro dos parametros saudaveis. Nenhum alerta de risco detectado.")
    else:
        for alert in alerts:
            color = (255, 79, 106) if alert["level"] == "error" else (255, 204, 79)
            pdf.set_text_color(*color)
            # Remove markdown bold markers para o PDF
            clean_msg = alert["msg"].replace("**", "").replace("*", "")
            prefix = "[CRITICO]" if alert["level"] == "error" else "[ATENCAO]"
            pdf.multi_cell(0, 6, f"{prefix} {alert['campaign']}: {clean_msg}")
            pdf.ln(1)
    pdf.ln(5)

    # ── Análise Gemini ──
    pdf.set_font("DejaVu", style="B", size=13)
    pdf.set_text_color(232, 230, 255)
    pdf.cell(0, 8, "Analise Estrategica — Gemini AI", ln=True)
    pdf.ln(2)

    pdf.set_font("DejaVu", size=9)
    pdf.set_text_color(200, 198, 232)
    # Remove emojis e formatação markdown pesada para compatibilidade PDF
    clean_analysis = analysis_text.replace("###", "").replace("##", "").replace("**", "")
    pdf.multi_cell(0, 5.5, clean_analysis)

    return bytes(pdf.output())


def generate_pptx(
    platform: str,
    summary: dict,
    alerts: list[dict],
    analysis_text: str,
    sym: str,
) -> bytes:
    """
    Gera uma apresentação PPTX com 3 slides:
      - Slide 1: Capa com plataforma e data
      - Slide 2: KPIs e alertas de Stop Loss
      - Slide 3: Análise estratégica do Gemini
    Retorna os bytes do PPTX prontos para st.download_button.
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Paleta de cores
    C_BG      = RGBColor(0x0a, 0x0a, 0x0f)   # #0a0a0f
    C_SURFACE = RGBColor(0x12, 0x12, 0x1a)   # #12121a
    C_ACCENT  = RGBColor(0x7c, 0x6a, 0xff)   # #7c6aff
    C_ACCENT2 = RGBColor(0xff, 0x6a, 0x9b)   # #ff6a9b
    C_TEXT    = RGBColor(0xe8, 0xe6, 0xff)   # #e8e6ff
    C_MUTED   = RGBColor(0x6b, 0x68, 0x80)   # #6b6880
    C_SUCCESS = RGBColor(0x4f, 0xff, 0xb0)   # #4fffb0
    C_ERROR   = RGBColor(0xff, 0x4f, 0x6a)   # #ff4f6a
    C_WARN    = RGBColor(0xff, 0xcc, 0x4f)   # #ffcc4f

    blank_layout = prs.slide_layouts[6]  # layout em branco

    def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        return shape

    def add_text_box(slide, text, l, t, w, h, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        txb.word_wrap = wrap
        tf = txb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or C_TEXT
        return txb

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 1 — Capa
    # ────────────────────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank_layout)
    add_rect(s1, 0, 0, 13.33, 7.5, C_BG)                      # fundo escuro
    add_rect(s1, 0, 0, 13.33, 0.12, C_ACCENT)                  # barra topo roxa
    add_rect(s1, 0, 7.38, 13.33, 0.12, C_ACCENT2)              # barra base rosa
    add_rect(s1, 5.5, 2.5, 2.33, 0.06, C_ACCENT)               # linha decorativa

    add_text_box(s1, "Marketing Digital Hub", 1, 1.6, 11.33, 1.2,
                 size=40, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(s1, f"Plataforma: {platform}", 1, 2.9, 11.33, 0.7,
                 size=20, bold=False, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_text_box(s1, "Relatório de Performance com Diagnóstico de Risco", 1, 3.55, 11.33, 0.6,
                 size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_text_box(s1, f"Gerado em {datetime.now().strftime('%d/%m/%Y às %H:%M')}", 1, 6.5, 11.33, 0.5,
                 size=10, color=C_MUTED, align=PP_ALIGN.CENTER)

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 2 — KPIs + Stop Loss
    # ────────────────────────────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank_layout)
    add_rect(s2, 0, 0, 13.33, 7.5, C_BG)
    add_rect(s2, 0, 0, 13.33, 0.08, C_ACCENT)

    add_text_box(s2, "KPIs & Diagnóstico de Risco", 0.4, 0.18, 12, 0.65,
                 size=22, bold=True, color=C_ACCENT)

    # ── KPI grid (2 colunas) ──
    kpi_lines = _kpi_lines(summary, sym)
    col_w, row_h = 5.8, 0.55
    for idx, line in enumerate(kpi_lines[:10]):   # max 10 KPIs
        col = idx % 2
        row = idx // 2
        bx = 0.4 + col * 6.5
        by = 1.05 + row * row_h
        add_rect(s2, bx, by, col_w, row_h - 0.06, C_SURFACE)
        add_text_box(s2, line, bx + 0.15, by + 0.05, col_w - 0.3, row_h - 0.12,
                     size=11, color=C_TEXT)

    kpi_rows = (min(len(kpi_lines), 10) + 1) // 2
    alert_y = 1.05 + kpi_rows * row_h + 0.2

    # ── Stop Loss ──
    add_text_box(s2, "🛑 Diagnóstico de Risco (Stop Loss)", 0.4, alert_y, 12, 0.5,
                 size=14, bold=True, color=C_ERROR if alerts else C_SUCCESS)
    alert_y += 0.5

    if not alerts:
        add_rect(s2, 0.4, alert_y, 12.4, 0.5, RGBColor(0x0d, 0x2b, 0x22))
        add_text_box(s2, "✅  Todas as campanhas estão dentro dos parâmetros saudáveis.", 0.55, alert_y + 0.05,
                     12, 0.4, size=11, color=C_SUCCESS)
    else:
        for alert in alerts[:4]:   # max 4 alertas no slide
            col_bg = RGBColor(0x2b, 0x0d, 0x12) if alert["level"] == "error" else RGBColor(0x2b, 0x24, 0x0d)
            col_tx = C_ERROR if alert["level"] == "error" else C_WARN
            clean = alert["msg"].replace("**", "").replace("*", "")
            add_rect(s2, 0.4, alert_y, 12.4, 0.52, col_bg)
            add_text_box(s2, clean[:130], 0.55, alert_y + 0.04, 12, 0.44,
                         size=9, color=col_tx)
            alert_y += 0.58

    # ────────────────────────────────────────────────────────────────────────
    # SLIDE 3 — Análise Estratégica Gemini
    # ────────────────────────────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank_layout)
    add_rect(s3, 0, 0, 13.33, 7.5, C_BG)
    add_rect(s3, 0, 0, 13.33, 0.08, C_ACCENT2)

    add_text_box(s3, "Análise Estratégica — Gemini AI", 0.4, 0.18, 12, 0.65,
                 size=22, bold=True, color=C_ACCENT2)
    add_text_box(s3, f"Plataforma: {platform}", 0.4, 0.78, 12, 0.35,
                 size=10, color=C_MUTED)

    # Limpa markdown para o PPTX e trunca para caber no slide
    clean = (analysis_text
             .replace("###", "")
             .replace("##", "")
             .replace("**", "")
             .replace("*", ""))
    max_chars = 2200
    if len(clean) > max_chars:
        clean = clean[:max_chars] + "\n\n[...] Texto completo disponível no download .txt"

    add_text_box(s3, clean, 0.4, 1.2, 12.53, 6.0,
                 size=9, color=C_TEXT, wrap=True)

    # ── Salva para bytes ──
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

# ── [fim REGRA 2 — funções de exportação] ────────────────────────────────────


# Chart theme
CHART_THEME = dict(
    paper_bgcolor="rgba(0,0,0,0)",
    plot_bgcolor="rgba(18,18,26,0.6)",
    font=dict(family="Syne, sans-serif", color="#e8e6ff", size=12),
    xaxis=dict(gridcolor="#2a2a3a", zerolinecolor="#2a2a3a", tickfont=dict(size=11)),
    yaxis=dict(gridcolor="#2a2a3a", zerolinecolor="#2a2a3a", tickfont=dict(size=11)),
    legend=dict(
        bgcolor="rgba(18,18,26,0.8)",
        bordercolor="#2a2a3a",
        borderwidth=1,
        font=dict(size=11),
    ),
    margin=dict(l=10, r=10, t=40, b=10),
    hovermode="x unified",
)

COLORS = {
    "spend":       "#7c6aff",
    "conversions": "#ff6a9b",
    "clicks":      "#4fffb0",
    "impressions": "#ffcc4f",
    "ctr":         "#ff9f4f",
    "roas":        "#4fc3ff",
    "cpc":         "#ff4f6a",
    "cpm":         "#b06aff",
}


def detect_date_col(df: pd.DataFrame):
    """Return the first column that looks like a date, or None."""
    candidates = [c for c in df.columns if any(k in c for k in ("date", "data", "dia", "day", "week", "period", "mes", "month"))]
    for col in candidates:
        try:
            parsed = pd.to_datetime(df[col], dayfirst=True, errors="coerce")
            if parsed.notna().sum() > len(df) * 0.5:
                return col, parsed
        except Exception:
            continue
    return None, None


def chart_spend_vs_conversions(df_t: pd.DataFrame, sym: str) -> go.Figure:
    fig = make_subplots(specs=[[{"secondary_y": True}]])

    fig.add_trace(go.Scatter(
        x=df_t["_date"], y=df_t["spend"],
        name=f"Investimento ({sym})",
        mode="lines+markers",
        line=dict(color=COLORS["spend"], width=2.5),
        marker=dict(size=5),
        fill="tozeroy",
        fillcolor="rgba(124,106,255,0.08)",
        hovertemplate=f"<b>Investimento</b>: {sym} %{{y:,.2f}}<extra></extra>",
    ), secondary_y=False)

    if "conversions" in df_t.columns:
        fig.add_trace(go.Scatter(
            x=df_t["_date"], y=df_t["conversions"],
            name="Conversões",
            mode="lines+markers",
            line=dict(color=COLORS["conversions"], width=2.5, dash="dot"),
            marker=dict(size=6, symbol="diamond"),
            hovertemplate="<b>Conversões</b>: %{y:,.0f}<extra></extra>",
        ), secondary_y=True)

    fig.update_layout(
        title=dict(text="📈 Investimento vs. Conversões ao longo do tempo", font=dict(size=15)),
        **CHART_THEME,
    )
    fig.update_yaxes(title_text=f"Investimento ({sym})", secondary_y=False, title_font=dict(color=COLORS["spend"]))
    fig.update_yaxes(title_text="Conversões", secondary_y=True, title_font=dict(color=COLORS["conversions"]))
    return fig


def chart_clicks_ctr(df_t: pd.DataFrame) -> go.Figure:
    fig = make_subplots(specs=[[{"secondary_y": True}]])

    fig.add_trace(go.Bar(
        x=df_t["_date"], y=df_t["clicks"] if "clicks" in df_t.columns else [],
        name="Cliques",
        marker_color=COLORS["clicks"],
        opacity=0.75,
        hovertemplate="<b>Cliques</b>: %{y:,.0f}<extra></extra>",
    ), secondary_y=False)

    if "ctr" in df_t.columns:
        fig.add_trace(go.Scatter(
            x=df_t["_date"], y=df_t["ctr"],
            name="CTR (%)",
            mode="lines+markers",
            line=dict(color=COLORS["ctr"], width=2.5),
            marker=dict(size=5),
            hovertemplate="<b>CTR</b>: %{y:.2f}%<extra></extra>",
        ), secondary_y=True)

    fig.update_layout(
        title=dict(text="🖱️ Cliques & CTR diário", font=dict(size=15)),
        barmode="overlay",
        **CHART_THEME,
    )
    fig.update_yaxes(title_text="Cliques", secondary_y=False, title_font=dict(color=COLORS["clicks"]))
    fig.update_yaxes(title_text="CTR (%)", secondary_y=True, title_font=dict(color=COLORS["ctr"]))
    return fig


def chart_roas_cpc(df_t: pd.DataFrame, sym: str) -> go.Figure:
    fig = make_subplots(specs=[[{"secondary_y": True}]])

    if "roas" in df_t.columns:
        fig.add_trace(go.Scatter(
            x=df_t["_date"], y=df_t["roas"],
            name="ROAS",
            mode="lines+markers",
            line=dict(color=COLORS["roas"], width=3),
            marker=dict(size=6),
            fill="tozeroy",
            fillcolor="rgba(79,195,255,0.06)",
            hovertemplate="<b>ROAS</b>: %{y:.2f}x<extra></extra>",
        ), secondary_y=False)

    if "cpc" in df_t.columns:
        fig.add_trace(go.Scatter(
            x=df_t["_date"], y=df_t["cpc"],
            name=f"CPC ({sym})",
            mode="lines+markers",
            line=dict(color=COLORS["cpc"], width=2, dash="dash"),
            marker=dict(size=5),
            hovertemplate=f"<b>CPC</b>: {sym} %{{y:.2f}}<extra></extra>",
        ), secondary_y=True)

    fig.update_layout(
        title=dict(text="🚀 ROAS vs. CPC ao longo do tempo", font=dict(size=15)),
        **CHART_THEME,
    )
    fig.update_yaxes(title_text="ROAS (x)", secondary_y=False, title_font=dict(color=COLORS["roas"]))
    fig.update_yaxes(title_text=f"CPC ({sym})", secondary_y=True, title_font=dict(color=COLORS["cpc"]))
    return fig


def chart_spend_by_campaign(df: pd.DataFrame, sym: str):
    if "campaign_name" not in df.columns or "spend" not in df.columns:
        return None
    grp = (
        df.groupby("campaign_name")["spend"]
        .sum()
        .reset_index()
        .sort_values("spend", ascending=True)
        .tail(12)
    )
    fig = go.Figure(go.Bar(
        x=grp["spend"],
        y=grp["campaign_name"],
        orientation="h",
        marker=dict(
            color=grp["spend"],
            colorscale=[[0, "#2a2a3a"], [0.5, "#7c6aff"], [1, "#ff6a9b"]],
            showscale=False,
        ),
        hovertemplate=f"<b>%{{y}}</b><br>Investimento: {sym} %{{x:,.2f}}<extra></extra>",
    ))
    
    try:
        fig.update_layout(**CHART_THEME)
    except:
        pass 
    
    fig.update_layout(
        title=dict(text="💸 Investimento por Campanha", font=dict(size=15)),
        yaxis=dict(tickfont=dict(size=10), gridcolor="#2a2a3a"),
        xaxis=dict(tickfont=dict(size=11), gridcolor="#2a2a3a"),
        height=int(max(300, len(grp) * 42)), 
    )
    return fig


def chart_funnel(df: pd.DataFrame) -> go.Figure:
    stages, values = [], []
    for col, label in [
        ("impressions", "Impressões"),
        ("reach",       "Alcance"),
        ("clicks",      "Cliques"),
        ("conversions", "Conversões"),
    ]:
        if col in df.columns:
            v = df[col].sum()
            if v > 0:
                stages.append(label)
                values.append(v)

    if len(stages) < 2:
        return None

    fig = go.Figure(go.Funnel(
        y=stages,
        x=values,
        textinfo="value+percent previous",
        marker=dict(color=["#7c6aff", "#a080ff", "#ff6a9b", "#4fffb0"]),
        connector=dict(line=dict(color="#2a2a3a", width=2)),
        hovertemplate="<b>%{y}</b>: %{x:,.0f}<extra></extra>",
    ))
    fig.update_layout(
        title=dict(text="🔽 Funil de Performance", font=dict(size=15)),
        **CHART_THEME,
        height=360,
    )
    return fig


# Sidebar
with st.sidebar:
    st.markdown("### ⚙️ Configuração")
    st.markdown("---")

    # Regra 1: st.sidebar.radio com as duas plataformas suportadas.
    # A variável `platform` controla o mapeamento de colunas, a limpeza do
    # CSV e o papel assumido pelo Gemini no prompt.
    platform = st.radio(
        "Escolha a Plataforma",
        options=["Meta Ads", "Google Ads"],
        horizontal=True,
        help="Selecione a origem do CSV para normalização automática das colunas.",
    )
    st.markdown("---")
 
    api_key = st.text_input(
        "🔑 Gemini API Key",
        type="password",
        placeholder="AIza...",
        help="Obtenha em aistudio.google.com",
    )

    model_choice = st.selectbox(
        "🤖 Modelo Gemini",
        ["gemini-1.5-pro", "gemini-2.0-flash", "gemini-2.5-flash"],
        index=1,
    )

    currency = st.selectbox(
        "💱 Moeda",
        ["BRL (R$)", "USD ($)", "EUR (€)", "GBP (£)"],
        index=0,
    )

    user_goal = st.text_area(
        "🎯 Objetivo da análise",
        placeholder="Ex: Reduzir CPA em campanhas de conversão de e-commerce...",
        height=110,
    )

    st.markdown("---")
    if platform == "Meta Ads":
        st.markdown("""
        <div style='font-family: Space Mono, monospace; font-size: 0.65rem; color: #6b6880; line-height: 1.8'>
        📋 <b>Colunas Meta Ads:</b><br>
        campaign_name · impressions<br>
        clicks · spend · reach<br>
        ctr · cpc · cpm · roas<br>
        conversions · frequency<br>
        ad_name · adset_name
        </div>
        """, unsafe_allow_html=True)
    else:
        st.markdown("""
        <div style='font-family: Space Mono, monospace; font-size: 0.65rem; color: #6b6880; line-height: 1.8'>
        📋 <b>Colunas Google Ads:</b><br>
        Campaign / Campanha<br>
        Impr. / Impressões<br>
        Clicks / Cliques<br>
        Cost / Custo<br>
        CTR · Avg. CPC / CPC médio<br>
        Conversions / Conversões<br>
        Conv. value / Valor conv.<br>
        Search Impr. share
        </div>
        """, unsafe_allow_html=True)


# Main
platform_icon = "🟦" if platform == "Meta Ads" else "🟥"
platform_label = platform

st.markdown(f"""
<div class="hero-header">
    <h1>{platform_icon} Marketing Digital Hub</h1>
    <p>Powered by Google Gemini · Analisando: <strong>{platform_label}</strong></p>
</div>
""", unsafe_allow_html=True)

upload_hint = (
    "Arraste ou selecione o CSV exportado do Meta Ads Manager"
    if platform == "Meta Ads"
    else "Arraste ou selecione o CSV exportado do Google Ads (Relatórios → Baixar → CSV)"
)

uploaded_file = st.file_uploader(
    upload_hint,
    type=["csv"],
    label_visibility="visible",
)

if uploaded_file:
    try:
        # Meta Ads usa o loader genérico original.
        if platform == "Google Ads":
            df = clean_google_ads_csv(uploaded_file)
        else:
            df = load_csv(uploaded_file)

        df = normalize_columns(df, platform)
        df = preprocess(df)

    except Exception as e:
        st.error(f"Erro ao carregar CSV: {e}")
        st.stop()

    # ── KPI Cards ──────────────────────────────────────────────────────────
    total_spend      = df["spend"].sum()          if "spend"            in df.columns else None
    total_impr       = df["impressions"].sum()    if "impressions"      in df.columns else None
    total_clicks     = df["clicks"].sum()         if "clicks"           in df.columns else None
    avg_ctr          = df["ctr"].mean()           if "ctr"              in df.columns else None
    avg_cpc          = df["cpc"].mean()           if "cpc"              in df.columns else None
    avg_roas         = df["roas"].mean()          if "roas"             in df.columns else None
    total_conv       = df["conversions"].sum()    if "conversions"      in df.columns else None
    total_revenue    = df["conversion_values"].sum() if "conversion_values" in df.columns else None

    sym = currency.split(" ")[1].strip("()")

    cards = []
    if total_spend    is not None: cards.append(("💸 Investimento",  fmt_currency(total_spend, sym), "", "purple"))
    if total_impr     is not None: cards.append(("👁️ Impressões",   fmt_number(total_impr),         "", "pink"))
    if total_clicks   is not None: cards.append(("🖱️ Cliques",      fmt_number(total_clicks),       "", "green"))
    if avg_ctr        is not None: cards.append(("📊 CTR Médio",    fmt_pct(avg_ctr),               "", "yellow"))
    if avg_cpc        is not None: cards.append(("💰 CPC Médio",    fmt_currency(avg_cpc, sym),     "", "purple"))
    if avg_roas       is not None: cards.append(("🚀 ROAS Médio",   f"{avg_roas:.2f}x",             "", "green"))
    if total_conv     is not None: cards.append(("✅ Conversões",   fmt_number(total_conv),          "", "pink"))
    if total_revenue  is not None: cards.append(("💵 Receita",      fmt_currency(total_revenue, sym),"", "yellow"))

    if cards:
        cols_html = "".join(
            f'<div class="kpi-card {color}">'
            f'<div class="kpi-label">{label}</div>'
            f'<div class="kpi-value">{value}</div>'
            f'<div class="kpi-sub">{sub}</div>'
            f'</div>'
            for label, value, sub, color in cards
        )
        st.markdown(f'<div class="kpi-grid">{cols_html}</div>', unsafe_allow_html=True)

    # ── Data preview ───────────────────────────────────────────────────────
    st.markdown('<div class="section-title">📋 Dados Importados</div>', unsafe_allow_html=True)
    st.dataframe(df, use_container_width=True, height=260)
    st.caption(f"{len(df)} linhas · {len(df.columns)} colunas detectadas")

    # ── Charts ─────────────────────────────────────────────────────────────
    st.markdown('<div class="section-title">📊 Visualizações</div>', unsafe_allow_html=True)

    date_col, date_series = detect_date_col(df)
    has_time = date_col is not None

    if has_time:
        df["_date"] = date_series
        agg_cols = {c: "sum" for c in ["spend", "clicks", "impressions", "conversions", "reach", "conversion_values"]
                    if c in df.columns}
        mean_cols = {c: "mean" for c in ["ctr", "cpc", "cpm", "roas", "frequency"] if c in df.columns}
        agg_cols.update(mean_cols)

        df_t = df.groupby("_date").agg(agg_cols).reset_index().sort_values("_date")

        # ── Row 1: Invest vs Conversions (full width) ──
        if "spend" in df_t.columns:
            fig1 = chart_spend_vs_conversions(df_t, sym)
            st.plotly_chart(fig1, use_container_width=True, config={"displayModeBar": False})

        # ── Row 2: Clicks+CTR | ROAS+CPC ──
        col_a, col_b = st.columns(2)
        with col_a:
            if "clicks" in df_t.columns or "ctr" in df_t.columns:
                st.plotly_chart(
                    chart_clicks_ctr(df_t),
                    use_container_width=True,
                    config={"displayModeBar": False},
                )
        with col_b:
            if "roas" in df_t.columns or "cpc" in df_t.columns:
                st.plotly_chart(
                    chart_roas_cpc(df_t, sym),
                    use_container_width=True,
                    config={"displayModeBar": False},
                )
    else:
        st.info(
            "💡 Nenhuma coluna de data detectada. Adicione uma coluna `date` ao CSV "
            "para ver os gráficos de linha temporais.",
            icon="📅",
        )

    # ── Row 3: Spend by campaign | Funnel ──
    col_c, col_d = st.columns([3, 2])
    with col_c:
        fig_camp = chart_spend_by_campaign(df, sym)
        if fig_camp:
            st.plotly_chart(fig_camp, use_container_width=True, config={"displayModeBar": False})
    with col_d:
        fig_funnel = chart_funnel(df)
        if fig_funnel:
            st.plotly_chart(fig_funnel, use_container_width=True, config={"displayModeBar": False})

    # ── [REGRA 1] Stop Loss — Diagnóstico de Risco ────────────────────────
    # Posicionado logo acima da seção de Análise com IA, conforme solicitado.
    # Os alertas são calculados aqui e também armazenados em session_state
    # para que a seção de exportação (mais abaixo) possa acessá-los.
    st.markdown('<div class="section-title">🛑 Gestão de Risco</div>', unsafe_allow_html=True)

    risk_alerts = run_stop_loss(df, sym)
    # Persiste alertas para uso no gerador de relatórios
    st.session_state["risk_alerts"] = risk_alerts

    with st.expander("🛑 Diagnóstico de Risco (Stop Loss)", expanded=bool(risk_alerts)):
        if not risk_alerts:
            st.success(
                "✅ **Todas as campanhas estão saudáveis!** "
                "Nenhum alerta de queima de verba, baixa relevância ou prejuízo detectado.",
                icon="✅",
            )
        else:
            # Contadores por severidade para o resumo no topo
            n_errors   = sum(1 for a in risk_alerts if a["level"] == "error")
            n_warnings = sum(1 for a in risk_alerts if a["level"] == "warning")
            st.markdown(
                f"**{len(risk_alerts)} alerta(s) encontrado(s):** "
                f"{n_errors} crítico(s) 🔴  ·  {n_warnings} atenção 🟡"
            )
            st.markdown("---")
            for alert in risk_alerts:
                camp_label = f"**Campanha:** `{alert['campaign']}`  —  "
                if alert["level"] == "error":
                    st.error(camp_label + alert["msg"])
                else:
                    st.warning(camp_label + alert["msg"])
    # ── [fim REGRA 1] ────────────────────────────────────────────────────────

    # ── Análise com IA ────────────────────────────────────────────────────────
    st.markdown('<div class="section-title">🤖 Análise com IA</div>', unsafe_allow_html=True)

    col_btn, col_info = st.columns([1, 3])
    with col_btn:
        run_analysis = st.button("✨ Analisar com Gemini", use_container_width=True)

    if run_analysis:
        if not api_key:
            st.warning("⚠️ Insira sua Gemini API Key na barra lateral.")
        else:
            summary = build_summary(df)
            goal    = user_goal or "Maximizar o desempenho geral das campanhas e reduzir custos."
            prompt  = build_prompt(summary, goal, currency, platform)

            with st.spinner("Gemini está analisando seus dados…"):
                try:
                    analysis = run_gemini_analysis(api_key, prompt, model_choice)

                    # [REGRA 2] Persiste o texto da análise para o gerador de relatórios
                    st.session_state["gemini_analysis"] = analysis
                    st.session_state["report_summary"]  = summary
                    st.session_state["report_sym"]      = sym

                    st.markdown("""
                    <div class="ai-box">
                        <span class="ai-tag">✦ Gemini AI · Análise de Performance</span>
                    """, unsafe_allow_html=True)
                    st.markdown(analysis)
                    st.markdown("</div>", unsafe_allow_html=True)

                    download_name = (
                        "meta_ads_analysis.txt"
                        if platform == "Meta Ads"
                        else "google_ads_analysis.txt"
                    )
                    st.download_button(
                        label="⬇️ Baixar análise (.txt)",
                        data=analysis.encode("utf-8"),
                        file_name=download_name,
                        mime="text/plain",
                    )
                except Exception as e:
                    st.error(f"Erro na API do Gemini: {e}")

    # ── [REGRA 2] Download de Relatórios (PDF e PPTX) ────────────────────────
    # A seção só fica ativa após a análise do Gemini ter sido executada ao
    # menos uma vez (dados armazenados em session_state).
    st.markdown('<div class="section-title">📥 Download de Relatórios</div>', unsafe_allow_html=True)

    has_report = (
        "gemini_analysis" in st.session_state
        and st.session_state["gemini_analysis"]
    )

    if not has_report:
        st.info(
            "💡 Execute a **Análise com Gemini** acima para liberar a exportação de relatórios.",
            icon="🔒",
        )
    else:
        st.markdown(
            "Relatórios gerados com os KPIs, alertas de risco e a análise completa do Gemini."
        )

        col_pdf, col_pptx, col_spacer = st.columns([1, 1, 2])

        # ── Botão PDF ──────────────────────────────────────────────────────
        with col_pdf:
            with st.spinner("Gerando PDF…"):
                try:
                    pdf_bytes = generate_pdf(
                        platform     = platform,
                        summary      = st.session_state["report_summary"],
                        alerts       = st.session_state.get("risk_alerts", []),
                        analysis_text= st.session_state["gemini_analysis"],
                        sym          = st.session_state["report_sym"],
                    )
                    st.download_button(
                        label     = "📄 Baixar PDF",
                        data      = pdf_bytes,
                        file_name = f"{platform.lower().replace(' ', '_')}_relatorio.pdf",
                        mime      = "application/pdf",
                        use_container_width=True,
                    )
                except Exception as e:
                    st.error(f"Erro ao gerar PDF: {e}")

        # ── Botão PPTX ─────────────────────────────────────────────────────
        with col_pptx:
            with st.spinner("Gerando PPTX…"):
                try:
                    pptx_bytes = generate_pptx(
                        platform     = platform,
                        summary      = st.session_state["report_summary"],
                        alerts       = st.session_state.get("risk_alerts", []),
                        analysis_text= st.session_state["gemini_analysis"],
                        sym          = st.session_state["report_sym"],
                    )
                    st.download_button(
                        label     = "📊 Baixar PPTX",
                        data      = pptx_bytes,
                        file_name = f"{platform.lower().replace(' ', '_')}_apresentacao.pptx",
                        mime      = "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True,
                    )
                except Exception as e:
                    st.error(f"Erro ao gerar PPTX: {e}")
    # ── [fim REGRA 2 — seção de exportação] ───────────────────────────────────

else:

    if platform == "Meta Ads":
        hint_platform = "Meta Ads Manager"
        hint_path     = "Relatórios → Exportar → Formato CSV"
        hint_cols     = "Campanha · Impressões · Cliques · Investimento · CTR"
    else:
        hint_platform = "Google Ads"
        hint_path     = "Relatórios → Baixar ícone → CSV"
        hint_cols     = "Campaign · Impr. · Clicks · Cost · CTR · Conversions"

    st.markdown(f"""
    <div style="
        text-align: center;
        padding: 4rem 2rem;
        color: #6b6880;
        font-family: 'Space Mono', monospace;
        font-size: 0.85rem;
        border: 1px dashed #2a2a3a;
        border-radius: 16px;
        margin-top: 2rem;
    ">
        <div style="font-size: 3rem; margin-bottom: 1rem;">📂</div>
        <div>Faça upload de um CSV exportado do</div>
        <div style="color: #7c6aff; font-weight: 700; margin-top: 0.3rem;">{hint_platform}</div>
        <div style="margin-top: 1rem; font-size: 0.72rem; line-height: 1.8; color: #4a4860">
            {hint_path}<br>
            Inclua: {hint_cols}
        </div>
    </div>
    """, unsafe_allow_html=True)